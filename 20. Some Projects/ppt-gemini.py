from typing import Dict, Any, List
from langgraph.graph import StateGraph, END
from langchain_core.messages import HumanMessage, AIMessage
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_core.output_parsers import JsonOutputParser
from pydantic import BaseModel, Field
import json
import re
import random
import requests
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os
from dotenv import load_dotenv

load_dotenv()


# Define Pydantic models for structured output
class PresentationOutline(BaseModel):
    title: str = Field(description="Presentation title")
    topic: str = Field(description="Main topic/subject")
    audience: str = Field(description="Target audience")
    tone: str = Field(
        description="Tone of presentation (Professional/Casual/Academic/etc)"
    )
    estimated_slides: int = Field(description="Number of slides (5-20)")
    key_points: List[str] = Field(description="Main points to cover")


class SlideContent(BaseModel):
    slide_number: int = Field(description="Slide number")
    title: str = Field(description="Slide title")
    content_type: str = Field(
        description="Type of slide: title_slide/content/bullet_points/conclusion/thank_you/image_slide"
    )
    main_content: str = Field(description="Main text content")
    bullet_points: List[str] = Field(
        description="Bullet points if applicable", default=[]
    )
    notes: str = Field(description="Speaker notes", default="")
    image_description: str = Field(
        description="Description for image search", default=""
    )


class SlidesData(BaseModel):
    slides: List[SlideContent] = Field(description="List of slides")


# Define the state structure
class WorkflowState(BaseModel):
    prompt: str = ""
    outline: Dict[str, Any] = {}
    slides_content: List[Dict[str, Any]] = []
    presentation_path: str = ""
    error: str = ""
    theme_name: str = ""  # Store theme name separately
    layout_style: str = ""  # Store layout style separately


# Design themes with different color schemes and layouts
DESIGN_THEMES = {
    "corporate_blue": {
        "primary_color": RGBColor(0, 123, 191),
        "secondary_color": RGBColor(51, 51, 51),
        "accent_color": RGBColor(0, 162, 232),
        "background_accent": RGBColor(240, 248, 255),
        "font": "Calibri",
        "accent_shape": MSO_SHAPE.RECTANGLE,
    },
    "modern_green": {
        "primary_color": RGBColor(46, 125, 50),
        "secondary_color": RGBColor(66, 66, 66),
        "accent_color": RGBColor(76, 175, 80),
        "background_accent": RGBColor(245, 255, 245),
        "font": "Segoe UI",
        "accent_shape": MSO_SHAPE.ROUNDED_RECTANGLE,
    },
    "elegant_purple": {
        "primary_color": RGBColor(103, 58, 183),
        "secondary_color": RGBColor(74, 74, 74),
        "accent_color": RGBColor(156, 39, 176),
        "background_accent": RGBColor(248, 245, 255),
        "font": "Arial",
        "accent_shape": MSO_SHAPE.HEXAGON,
    },
    "warm_orange": {
        "primary_color": RGBColor(255, 87, 34),
        "secondary_color": RGBColor(62, 62, 62),
        "accent_color": RGBColor(255, 152, 0),
        "background_accent": RGBColor(255, 248, 240),
        "font": "Tahoma",
        "accent_shape": MSO_SHAPE.OVAL,
    },
    "tech_teal": {
        "primary_color": RGBColor(0, 150, 136),
        "secondary_color": RGBColor(55, 55, 55),
        "accent_color": RGBColor(0, 188, 212),
        "background_accent": RGBColor(240, 255, 255),
        "font": "Verdana",
        "accent_shape": MSO_SHAPE.DIAMOND,
    },
}

# Layout variations
LAYOUT_STYLES = [
    "clean_minimal",
    "sidebar_accent",
    "header_focus",
    "centered_content",
    "asymmetric_modern",
]


# Initialize the LLM
llm = ChatGoogleGenerativeAI(model="gemini-2.5-flash", temperature=0.7)

# Create output parsers
outline_parser = JsonOutputParser(pydantic_object=PresentationOutline)
slides_parser = JsonOutputParser(pydantic_object=SlidesData)


def get_unsplash_image(query, width=800, height=600):
    """
    Get a free image from Unsplash API
    """
    try:
        # Unsplash API endpoint for random photos
        url = f"https://source.unsplash.com/{width}x{height}/?{query}"

        response = requests.get(url, timeout=10)
        if response.status_code == 200:
            return BytesIO(response.content)
        else:
            print(f"Could not fetch image for: {query}")
            return None
    except Exception as e:
        print(f"Error fetching image: {e}")
        return None


def select_random_theme():
    """
    Select a random design theme for variety
    """
    theme_name = random.choice(list(DESIGN_THEMES.keys()))
    layout_style = random.choice(LAYOUT_STYLES)

    print(f"Selected theme: {theme_name} with {layout_style} layout")
    return theme_name, layout_style


def analyze_prompt_node(state: WorkflowState) -> WorkflowState:
    """
    Analyze the user prompt and extract key information about the presentation
    """
    try:
        # Select random theme for this presentation
        theme_name, layout_style = select_random_theme()
        state.theme_name = theme_name
        state.layout_style = layout_style

        analysis_prompt = f"""
        Analyze this presentation request and extract key information:
        
        User Request: {state.prompt}
        
        {outline_parser.get_format_instructions()}
        """

        # Create a chain with the parser
        chain = llm | outline_parser
        analysis = chain.invoke([HumanMessage(content=analysis_prompt)])

        # Convert to dict for state storage
        state.outline = analysis if isinstance(analysis, dict) else analysis.dict()
        print(f"Analysis completed: {state.outline.get('title', 'Unknown')}")
        return state

    except Exception as e:
        print(f"Error in analyze_prompt_node: {str(e)}")
        state.error = f"Error analyzing prompt: {str(e)}"
        return state


def generate_content_node(state: WorkflowState) -> WorkflowState:
    """
    Generate detailed content for each slide based on the outline with varied formats
    """
    try:
        if state.error:
            return state

        estimated_slides = state.outline.get("estimated_slides", 8)
        total_slides = estimated_slides + 1

        content_prompt = f"""
        Based on this presentation outline, generate detailed content for {total_slides} slides with VARIED FORMATS.
        
        Outline: {json.dumps(state.outline, indent=2)}
        
        VARIETY REQUIREMENTS:
        - Mix different content_types: "bullet_points", "content", "image_slide"
        - Include 2-3 "image_slide" slides with relevant image_description
        - Vary slide structures and layouts
        - Keep content concise (bullet points: max 8 words, max 5 points)
        - Make it visually interesting and dynamic
        
        Content Types to Use:
        - title_slide: Opening slide
        - bullet_points: Key points in bullets
        - content: Paragraph text
        - image_slide: Slide with image and minimal text
        - conclusion: Summary slide
        - thank_you: Final slide
        
        For image_slide slides, provide descriptive image_description for image search.
        
        Example image_description: "renewable energy solar panels business office", "team meeting corporate environment", "growth chart business success"
        
        {slides_parser.get_format_instructions()}
        """

        chain = llm | slides_parser
        content_data = chain.invoke([HumanMessage(content=content_prompt)])

        # Convert to dict format
        if isinstance(content_data, dict):
            slides = content_data.get("slides", [])
        else:
            slides = content_data.slides

        slides_dict = []
        for slide in slides:
            if isinstance(slide, dict):
                slides_dict.append(slide)
            else:
                slides_dict.append(slide.dict())

        # Ensure we have a thank you slide
        has_thank_you = any(
            slide.get("content_type") == "thank_you" for slide in slides_dict
        )
        if not has_thank_you:
            thank_you_slide = {
                "slide_number": len(slides_dict) + 1,
                "title": "Thank You",
                "content_type": "thank_you",
                "main_content": "Questions & Discussion",
                "bullet_points": [],
                "notes": "Thank the audience and open for questions",
                "image_description": "",
            }
            slides_dict.append(thank_you_slide)

        state.slides_content = slides_dict
        print(f"Generated {len(slides_dict)} slides with varied formats")
        return state

    except Exception as e:
        print(f"Error in generate_content_node: {str(e)}")
        state.error = f"Error generating content: {str(e)}"
        return state


def apply_theme_styling(slide, theme_colors, layout_style, slide_data):
    """
    Apply dynamic theme styling based on selected theme and layout
    """
    try:
        if layout_style == "sidebar_accent":
            # Wider accent bar on the side
            accent_width = Inches(0.5)
        elif layout_style == "clean_minimal":
            # Thin accent line
            accent_width = Inches(0.2)
        else:
            # Standard accent
            accent_width = Inches(0.3)

        # Add themed accent shape
        if layout_style == "header_focus":
            # Top accent bar
            left, top, width, height = Inches(0), Inches(0), Inches(13.33), Inches(0.3)
        elif layout_style == "sidebar_accent":
            # Side accent bar
            left, top, width, height = Inches(0), Inches(0), accent_width, Inches(7.5)
        else:
            # Standard left accent
            left, top, width, height = Inches(0), Inches(0), accent_width, Inches(7.5)

        accent_shape = slide.shapes.add_shape(
            theme_colors["accent_shape"], left, top, width, height
        )

        fill = accent_shape.fill
        fill.solid()
        fill.fore_color.rgb = theme_colors["primary_color"]
        accent_shape.line.fill.background()

    except Exception as e:
        print(f"Warning: Could not apply theme styling: {e}")


def create_image_slide(slide, slide_data, theme_colors):
    """
    Create a slide with image and text overlay
    """
    try:
        # Get image if description provided
        if slide_data.get("image_description"):
            image_data = get_unsplash_image(slide_data["image_description"], 800, 600)

            if image_data:
                # Add image
                left = Inches(1)
                top = Inches(1.5)
                width = Inches(11.33)
                height = Inches(5)

                pic = slide.shapes.add_picture(image_data, left, top, width, height)

                # Add semi-transparent overlay for text readability
                overlay = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, left, top + Inches(3), width, Inches(2)
                )
                fill = overlay.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(0, 0, 0)
                overlay.fill.transparency = 0.3
                overlay.line.fill.background()

                # Add text over image
                textbox = slide.shapes.add_textbox(
                    left + Inches(0.5),
                    top + Inches(3.2),
                    width - Inches(1),
                    Inches(1.5),
                )
                text_frame = textbox.text_frame
                text_frame.text = slide_data.get("main_content", "")

                # Style text
                for paragraph in text_frame.paragraphs:
                    paragraph.font.size = Pt(24)
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = RGBColor(255, 255, 255)
                    paragraph.alignment = PP_ALIGN.CENTER

        # Add title
        title_shape = slide.shapes.add_textbox(
            Inches(1), Inches(0.5), Inches(11.33), Inches(1)
        )
        title_frame = title_shape.text_frame
        title_frame.text = slide_data.get("title", "")

        for paragraph in title_frame.paragraphs:
            paragraph.font.size = Pt(36)
            paragraph.font.bold = True
            paragraph.font.color.rgb = theme_colors["primary_color"]
            paragraph.alignment = PP_ALIGN.LEFT

    except Exception as e:
        print(f"Warning: Could not create image slide: {e}")
        # Fallback to regular slide
        create_regular_slide(slide, slide_data, theme_colors, "centered_content")


def create_regular_slide(slide, slide_data, theme_colors, layout_style):
    """
    Create regular content slides with varied layouts
    """
    try:
        if layout_style == "centered_content":
            # Centered layout
            title_left, title_top = Inches(1), Inches(1)
            content_left, content_top = Inches(2), Inches(2.5)
            content_width = Inches(9.33)
            text_align = PP_ALIGN.CENTER
        elif layout_style == "asymmetric_modern":
            # Asymmetric layout
            title_left, title_top = Inches(1.5), Inches(0.8)
            content_left, content_top = Inches(1.5), Inches(2.2)
            content_width = Inches(10)
            text_align = PP_ALIGN.LEFT
        else:
            # Standard layout
            title_left, title_top = Inches(0.7), Inches(0.5)
            content_left, content_top = Inches(0.7), Inches(1.8)
            content_width = Inches(11.63)
            text_align = PP_ALIGN.LEFT

        # Add title
        title_shape = slide.shapes.add_textbox(
            title_left, title_top, Inches(11), Inches(1)
        )
        title_frame = title_shape.text_frame
        title_frame.text = slide_data.get("title", "")

        for paragraph in title_frame.paragraphs:
            paragraph.font.size = Pt(36)
            paragraph.font.bold = True
            paragraph.font.color.rgb = theme_colors["primary_color"]
            paragraph.font.name = theme_colors["font"]
            paragraph.alignment = text_align

        # Add content
        content_shape = slide.shapes.add_textbox(
            content_left, content_top, content_width, Inches(4.5)
        )
        content_frame = content_shape.text_frame
        content_frame.margin_left = Inches(0.3)
        content_frame.margin_top = Inches(0.2)

        if slide_data.get("content_type") == "bullet_points" and slide_data.get(
            "bullet_points"
        ):
            for j, bullet in enumerate(slide_data["bullet_points"]):
                p = (
                    content_frame.paragraphs[0]
                    if j == 0
                    else content_frame.add_paragraph()
                )
                p.text = str(bullet)
                p.level = 0
                p.font.size = Pt(24)
                p.font.name = theme_colors["font"]
                p.font.color.rgb = theme_colors["secondary_color"]
                p.space_after = Pt(16)
                p.alignment = text_align
        else:
            content_frame.text = slide_data.get("main_content", "")
            for paragraph in content_frame.paragraphs:
                paragraph.font.size = Pt(20)
                paragraph.font.name = theme_colors["font"]
                paragraph.font.color.rgb = theme_colors["secondary_color"]
                paragraph.space_after = Pt(12)
                paragraph.alignment = text_align

    except Exception as e:
        print(f"Warning: Could not create regular slide: {e}")


def create_presentation_node(state: WorkflowState) -> WorkflowState:
    """
    Create presentation with dynamic themes and image support
    """
    try:
        if state.error:
            return state

        # Get theme and layout from state (fixed version)
        theme_name = state.theme_name
        layout_style = state.layout_style
        theme_colors = DESIGN_THEMES[theme_name]

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        print(
            f"Creating presentation with {len(state.slides_content)} slides using {theme_name} theme and {layout_style} layout"
        )

        for i, slide_data in enumerate(state.slides_content):
            content_type = slide_data.get("content_type", "content")
            print(
                f"Processing slide {i+1}: {slide_data.get('title', 'Untitled')} - {content_type}"
            )

            # Use blank layout for custom positioning
            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)

            # Apply theme styling
            apply_theme_styling(slide, theme_colors, layout_style, slide_data)

            if content_type == "title_slide":
                # Title slide with theme colors
                title_shape = slide.shapes.add_textbox(
                    Inches(1), Inches(2.5), Inches(11.33), Inches(1.5)
                )
                title_frame = title_shape.text_frame
                title_frame.text = slide_data.get("title", "")

                for paragraph in title_frame.paragraphs:
                    paragraph.font.size = Pt(48)
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = theme_colors["primary_color"]
                    paragraph.font.name = theme_colors["font"]
                    paragraph.alignment = PP_ALIGN.CENTER

                if slide_data.get("main_content"):
                    subtitle_shape = slide.shapes.add_textbox(
                        Inches(1), Inches(4.5), Inches(11.33), Inches(1)
                    )
                    subtitle_frame = subtitle_shape.text_frame
                    subtitle_frame.text = slide_data["main_content"]

                    for paragraph in subtitle_frame.paragraphs:
                        paragraph.font.size = Pt(28)
                        paragraph.font.color.rgb = theme_colors["secondary_color"]
                        paragraph.font.name = theme_colors["font"]
                        paragraph.alignment = PP_ALIGN.CENTER

            elif content_type == "image_slide":
                create_image_slide(slide, slide_data, theme_colors)

            elif content_type == "thank_you":
                # Styled thank you slide
                title_shape = slide.shapes.add_textbox(
                    Inches(1), Inches(2), Inches(11.33), Inches(2)
                )
                title_frame = title_shape.text_frame
                title_frame.text = slide_data.get("title", "Thank You")

                for paragraph in title_frame.paragraphs:
                    paragraph.font.size = Pt(54)
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = theme_colors["primary_color"]
                    paragraph.font.name = theme_colors["font"]
                    paragraph.alignment = PP_ALIGN.CENTER

                if slide_data.get("main_content"):
                    content_shape = slide.shapes.add_textbox(
                        Inches(1), Inches(4.5), Inches(11.33), Inches(1.5)
                    )
                    content_frame = content_shape.text_frame
                    content_frame.text = slide_data["main_content"]

                    for paragraph in content_frame.paragraphs:
                        paragraph.font.size = Pt(32)
                        paragraph.font.color.rgb = theme_colors["secondary_color"]
                        paragraph.font.name = theme_colors["font"]
                        paragraph.alignment = PP_ALIGN.CENTER

            else:
                create_regular_slide(slide, slide_data, theme_colors, layout_style)

            # Add speaker notes
            if slide_data.get("notes"):
                try:
                    notes_slide = slide.notes_slide
                    notes_text_frame = notes_slide.notes_text_frame
                    notes_text_frame.text = slide_data["notes"]
                except Exception:
                    pass

        # Save with theme name in filename
        title = state.outline.get("title", "Generated_Presentation")
        filename = re.sub(r'[<>:"/\\|?*]', "_", title).replace(" ", "_")
        filename += f"_{theme_name}_{layout_style}.pptx"
        filepath = os.path.join(os.getcwd(), filename)
        prs.save(filepath)

        state.presentation_path = filepath
        print(f"Presentation saved to: {filepath}")
        return state

    except Exception as e:
        print(f"Error in create_presentation_node: {str(e)}")
        state.error = f"Error creating presentation: {str(e)}"
        return state


def error_handler_node(state: WorkflowState) -> WorkflowState:
    """Handle errors"""
    print(f"Error occurred: {state.error}")
    return state


def success_node(state: WorkflowState) -> WorkflowState:
    """Final success node"""
    print(f"Presentation successfully created: {state.presentation_path}")
    return state


# Define routing logic
def route_after_analysis(state: WorkflowState) -> str:
    return "error_handler" if state.error else "generate_content"


def route_after_content(state: WorkflowState) -> str:
    return "error_handler" if state.error else "create_presentation"


def route_after_creation(state: WorkflowState) -> str:
    return "error_handler" if state.error else "success"


# Create workflow
def create_ppt_workflow():
    workflow = StateGraph(WorkflowState)

    workflow.add_node("analyze_prompt", analyze_prompt_node)
    workflow.add_node("generate_content", generate_content_node)
    workflow.add_node("create_presentation", create_presentation_node)
    workflow.add_node("error_handler", error_handler_node)
    workflow.add_node("success", success_node)

    workflow.set_entry_point("analyze_prompt")

    workflow.add_conditional_edges(
        "analyze_prompt",
        route_after_analysis,
        {"generate_content": "generate_content", "error_handler": "error_handler"},
    )
    workflow.add_conditional_edges(
        "generate_content",
        route_after_content,
        {
            "create_presentation": "create_presentation",
            "error_handler": "error_handler",
        },
    )
    workflow.add_conditional_edges(
        "create_presentation",
        route_after_creation,
        {"success": "success", "error_handler": "error_handler"},
    )

    workflow.add_edge("error_handler", END)
    workflow.add_edge("success", END)

    return workflow.compile()


def generate_presentation(prompt: str):
    """Main function to generate a presentation"""
    app = create_ppt_workflow()
    initial_state = WorkflowState(prompt=prompt)
    result = app.invoke(initial_state)
    return result


if __name__ == "__main__":
    user_prompt = """
    Create a presentation about the benefits of renewable energy for a business audience. 
    Include information about cost savings, environmental impact, and implementation strategies.
    Make it professional and data-driven with about 10 slides.
    """

    result = generate_presentation(user_prompt)

    if result.get("error"):
        print(f"Failed to generate presentation: {result.get('error')}")
    else:
        print(f"Successfully generated presentation: {result.get('presentation_path')}")
        print(f"Title: {result.get('outline', {}).get('title')}")
        print(f"Number of slides: {len(result.get('slides_content', []))}")
        print(
            f"Theme used: {result.get('theme_name')} with {result.get('layout_style')} layout"
        )
